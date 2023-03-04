from flask import Flask, render_template
from pptx import Presentation
import spacy

app = Flask(__name__)

@app.route('/')
def index():
    # Load the PowerPoint file
    prs = Presentation('example.pptx')

    # Extract text from each slide
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text.append(shape.text)

    # Use spaCy to extract key phrases and generate summary sentences
    nlp = spacy.load('en_core_web_sm')
    doc = nlp('\n'.join(text))
    phrases = [chunk.text for chunk in doc.noun_chunks]
    sentences = [sent.text for sent in doc.sents]
    summary = '\n'.join(['- ' + sentence for sentence in sentences])

    # Render the summary on a web page
    return render_template('index.html', summary=summary)

if __name__ == '__main__':
    app.run(debug=True)
        